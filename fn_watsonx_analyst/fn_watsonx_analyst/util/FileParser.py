"""File parser module to parse Artifacts and Attachments in the incident"""

# pylint: disable=wildcard-import, method-hidden, import-outside-toplevel, import-error, broad-exception-caught, invalid-name

import os
import io
import re
import json
from io import BytesIO
import base64
from typing import Dict, Callable, Set, Any
from pytesseract import image_to_string
from bs4 import BeautifulSoup
from mailparser import parse_from_bytes
from xlrd import open_workbook
from pdf2image import convert_from_bytes
from pypdf import PdfReader
from docx import Document
from PIL import Image
from ics import Calendar
from pptx import Presentation
from fn_watsonx_analyst.util.logging_helper import create_logger

# Configure logging
log = create_logger(__name__)

class UnsupportedEncodingError(Exception):
    """Raised when a file is not encoded in UTF-8 or ASCII."""
    pass

class FileParser:
    """Class to parse multiple file types"""
    parsers: Dict[str, Callable[[Any], str]]
    supported_extensions: Set[str]
    
    def __init__(self):
        """Initialize the file parser with supported extensions and parser mapping."""
        self.parsers = {
            ".html": self.read_html,
            ".ics": self.read_ics,
            ".pdf": self.read_pdf,
            ".docx": self.read_docx,
            ".xlsx": self.read_xlsx,
            ".ipynb": self.read_ipynb,
            ".pptx": self.read_pptx,
            ".ppt": self.read_ppt,
            ".json": self.read_json,
            ".eml": self.read_eml,
            ".jpg": self.extract_text_from_images,
            ".jpeg": self.extract_text_from_images,
            ".png": self.extract_text_from_images,
        }
        self.supported_extensions = set(self.parsers.keys())

    def get_supported_extensions(self)-> Set[str]:
        """for better Abstraction"""
        log.info("Supported extensions requested")
        return self.supported_extensions

    def extract_text_from_images(self, images: list) -> str:
        """Extract text from a list of image bytes or PIL Images using OCR."""
        try:
            text = []
            if isinstance(images, bytes):
                images = [images]
            for i, img in enumerate(images):
                if isinstance(img, bytes):
                    try:
                        img = Image.open(io.BytesIO(img))
                    except Exception:
                        log.exception("Could not convert bytes to image at index: %s", i)
                        raise
                elif not isinstance(img, Image.Image):
                    raise TypeError(f"Unsupported type at index {i}: {type(img)}")

                ocr_text = image_to_string(img)
                if ocr_text.strip():
                    text.append(ocr_text)
            return self.clean_text("\n\n".join(text)) if text else ""
        except Exception as e:
            log.exception("OCR failed while converting images to text: %s", e)
            raise

    def clean_text(self, data: str) -> str:
        """Remove multiple empty lines and HTML tags."""
        try:
            text = data
            # Only attempt to parse HTML if there are likely HTML tags
            if re.search(r"<[^>]+>", data):
                # Parse the HTML data
                soup = BeautifulSoup(data, "html.parser")

                # Remove all HTML tags
                text = soup.get_text()

            # Remove multiple empty lines
            return re.sub(r"\n\s*\n+", "\n\n", text).strip()
        except Exception as e:
            log.exception("Error cleaning text %s: %s", data, e)
            raise

    def read_txt(self, data: str) -> str:
        """Read text files."""
        try:
            lines = []
            # Split the input string into lines
            for line in data.splitlines():
                stripped = line.strip()
                if stripped:  # Ignore empty lines
                    lines.append(line)
            return self.clean_text("\n".join(lines))
        except Exception as e:
            log.exception("Error reading text %s: %s", data, e)
            raise

    def read_json(self, data: str):
        """Read text files."""
        try:
            return self.clean_text(data)
        except Exception as e:
            log.exception("Error reading TXT stream: %s", e)
            raise

    def read_html(self, data: str) -> str:
        """Extract text & image-text from HTML (string only)."""
        try:
            if not isinstance(data, str):
                raise TypeError("Invalid HTML input: Expected a string.")

            soup = BeautifulSoup(data, "html.parser")
            text = soup.get_text(separator="\n")

            # Extract base64 images if present
            images = []
            for img_tag in soup.find_all("img"):
                src = img_tag.get("src")
                if src and "base64," in src:
                    try:
                        img_data = base64.b64decode(src.split("base64,")[1])
                        images.append(Image.open(BytesIO(img_data)))
                    except Exception as e:
                        log.warning("Skipping invalid base64 image: %s", e)

            ocr_text = self.extract_text_from_images(images) if images else ""
            full_text = "\n\n".join([text, ocr_text]) if ocr_text else text

            return (
                self.clean_text(full_text)
                if full_text
                else "No readable content in HTML."
            )

        except Exception as e:
            log.exception("Error reading HTML: %s", e)
            raise

    def read_ics(self, data: str) -> str:
        """Extract text from ICS calendar files."""
        try:
            if isinstance(data, str):
                data = data.encode("utf-8")  # Convert str to bytes if necessary

            decoded_data = data.decode("utf-8")

            # If PRODID is missing, add a default one else we get exception
            if "PRODID" not in decoded_data:
                decoded_data = "BEGIN:VCALENDAR\nPRODID:-//Default//EN\n" + decoded_data

            calendar = Calendar(
                data.decode("utf-8")
            )  # ICS data is often UTF-8 encoded text
            text = "\n".join(
                [
                    f"Event: {event.name}, Start: {event.begin}, End: {event.end}"
                    for event in calendar.events
                ]
            )
            return self.clean_text(text)
        except Exception as e:
            log.exception("Error reading ICS file: %s", e)
            raise

    def read_eml(self, data: str) -> str:
        """Extract plain text from EML data (bytes or str)."""
        try:
            if isinstance(data, str):
                data = data.encode("utf-8")  # Convert str to bytes
            mail = parse_from_bytes(data)
            text = mail.text_plain[0] if mail.text_plain else ""
            return text.strip()
        except Exception as e:
            print(f"Error reading EML: {e}")
            raise

    def read_xls(self, data: str) -> str:
        """Parses excel(.xls) file"""
        try:
            book = open_workbook(data)
            text = []
            for sheet in book.sheets():
                for row_idx in range(sheet.nrows):
                    text.append(
                        ", ".join(
                            str(sheet.cell(row_idx, col_idx).value)
                            for col_idx in range(sheet.ncols)
                        )
                    )
            return self.clean_text("\n".join(text))

        except ValueError as ve:
            log.error("excel(xls) artifact parsing error: %s", ve)
            raise

        except Exception as e:
            log.exception("An error occurred during parsing an excel(xls) file: %s", e)
            raise

    def read_xlsx(self, data: str) -> str:
        """Parses excel (.xlsx) file and extracts text."""
        try:
            if isinstance(data, str):
                data = data.encode("utf-8")  # Convert string to bytes if necessary
            from openpyxl import load_workbook

            # Load the workbook from the byte data
            wb = load_workbook(filename=io.BytesIO(data), read_only=True)
            text = []

            # Iterate over all sheets in the workbook
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                # Iterate over all rows in the sheet and extract values
                for row in ws.iter_rows(values_only=True):
                    text.append(", ".join(str(cell) if cell else "" for cell in row))

            return self.clean_text(
                "\n".join(text)
            )  # Clean and return the extracted text
        except Exception as e:
            log.exception("Error reading XLSX file: %s", e)
            raise

    def read_pdf(self, data: bytes) -> str:
        """Extract text and image-text from PDFs using PyPDF."""
        try:
            if not isinstance(data, bytes):
                return "Invalid PDF input: Expected bytes."

            text = []
            reader = PdfReader(BytesIO(data))

            for i, page in enumerate(reader.pages):
                try:
                    page_text = page.extract_text()
                    if page_text and page_text.strip():
                        text.append(page_text)
                    else:
                        log.debug(f"No text found on page %s: {i}" )
                except Exception as e:
                    log.exception("Failed to extract text from pdf file %s", e)

            if not text:
                # Fallback to OCR
                log.info("No extractable text found. Falling back to OCR.")
                images = convert_from_bytes(data)
                ocr_text = self.extract_text_from_images(images)
                if ocr_text:
                    text.append(ocr_text)

            return (
                self.clean_text("\n\n".join(text))
                if text
                else "No readable content found in PDF."
            )

        except Exception as e:
            log.exception("Error reading PDF: %s", e)
            raise

    def read_docx(self, data: str) -> str:
        """Extract text & image-text from DOCX (bytes only)."""
        try:
            if not isinstance(data, bytes):
                return "Invalid DOCX input: Expected bytes."

            doc = Document(BytesIO(data))
            text = [para.text for para in doc.paragraphs if para.text.strip()]

            # Extract images
            images = []
            for rel in doc.part._rels.values():
                if "image" in rel.target_ref:
                    image_bytes = rel.target_part.blob
                    images.append(Image.open(BytesIO(image_bytes)))

            ocr_text = self.extract_text_from_images(images) if images else ""
            full_text = (
                "\n\n".join(text + [ocr_text]) if ocr_text else "\n\n".join(text)
            )

            return (
                self.clean_text(full_text)
                if full_text
                else "No readable content in DOCX."
            )

        except Exception as e:
            log.exception("Error reading DOCX: %s", e)
            raise

    def read_pptx(self, data: bytes) -> str:
        """Extract text from PowerPoint (.pptx)."""
        try:
            if isinstance(data, str):
                data = data.encode("utf-8")  # Convert str to bytes if necessary

            prs = Presentation(BytesIO(data))
            text = []
            for slide in prs.slides:
                slide_text = [
                    shape.text.strip()
                    for shape in slide.shapes
                    if hasattr(shape, "text")
                ]
                if slide_text:
                    text.append("\n".join(slide_text))
            return self.clean_text("\n\n".join(text))
        except Exception as e:
            log.exception("Error reading PPTX file: %s", e)
            raise

    def convert_ppt_to_pptx(self, data: bytes) -> bytes:
        """Convert .ppt to .pptx using LibreOffice headless mode and return the converted file as bytes."""
        try:
            import subprocess
            import tempfile

            with tempfile.NamedTemporaryFile(suffix=".ppt", delete=False) as temp_ppt:
                temp_ppt.write(data)
                temp_ppt_path = temp_ppt.name

            output_pptx_path = temp_ppt_path.replace(".ppt", ".pptx")

            # Convert using LibreOffice (ensure it's installed)
            subprocess.run(
                [
                    "soffice",
                    "--headless",
                    "--convert-to",
                    "pptx",
                    temp_ppt_path,
                    "--outdir",
                    os.path.dirname(temp_ppt_path),
                ],
                check=True,
            )

            # Read converted .pptx file as bytes
            with open(output_pptx_path, "rb") as f:
                pptx_bytes = f.read()

            # Cleanup temp files
            os.remove(temp_ppt_path)
            os.remove(output_pptx_path)

            return pptx_bytes

        except Exception as e:
            log.exception("Error converting PPT to PPTX: %s", e)
            raise

    def read_ppt(self, data: bytes) -> str:
        """Convert .ppt to .pptx and extract text."""
        try:
            # import takes about 25MB so leaving it here for now
            from pptx import Presentation
            import subprocess

            # Convert .ppt to .pptx using unoconv (Linux/macOS only)
            pptx_data = self.convert_ppt_to_pptx(data)
            if pptx_data is None:
                return "Error: Cannot read .ppt file."

            # Extract text from .pptx
            prs = Presentation(io.BytesIO(pptx_data))
            text = []
            for slide in prs.slides:
                slide_text = [
                    shape.text.strip()
                    for shape in slide.shapes
                    if hasattr(shape, "text")
                ]
                if slide_text:
                    text.append("\n".join(slide_text))

            return self.clean_text("\n\n".join(text))

        except Exception as e:
            log.exception("Error reading PPT file: %s", e)
            raise

    def read_ipynb(self, data: str) -> str:
        """Parsing text from Jupyter Notebook (`.ipynb`) and remove HTML tags."""
        try:
            notebook = json.loads(data)
            text = []
            for cell in notebook.get("cells", []):
                if cell.get("cell_type") in ["code", "markdown"]:
                    raw_text = "\n".join(cell.get("source", []))
                    cleaned_text = self.clean_text(raw_text)  # Remove HTML tags
                    text.append(cleaned_text)

            return self.clean_text("\n".join(text))
        except Exception as e:
            log.exception(
                "An error occurred during Jupyter Notebook file parsing: %s", e
            )
            raise

    def multi_format_parser(self, data: any, object_name: str) -> str:
        """Determine file type and extract text."""
        ext = os.path.splitext(object_name)[-1].lower()
        parsed_content = ""

        try:
            if ext in self.get_supported_extensions():
                parser = self.parsers[ext]
                parsed_content = parser(data)
            else:
                # Try to process as raw UTF-8 or ASCII text
                if isinstance(data, bytes):
                    try:
                        decoded = data.decode('utf-8')
                    except UnicodeDecodeError:
                        try:
                            decoded = data.decode('ascii')
                        except UnicodeDecodeError as ue:
                            raise UnsupportedEncodingError("File is not UTF-8 or ASCII encoded.") from ue
                    parsed_content = self.read_txt(decoded)
                elif isinstance(data, str):
                    try:
                        data.encode('ascii')  # Check if string can be represented in ASCII
                    except UnicodeEncodeError as ue:
                        raise UnsupportedEncodingError("String data is not ASCII encodable.") from ue
                    parsed_content = self.read_txt(data)
                else:
                    raise UnsupportedEncodingError("Unsupported data type for decoding.")

            if not parsed_content or not parsed_content.strip():
                raise ValueError("Parsed content is empty or could not be extracted.")

        except Exception:
            parsed_content = "Parsed content is empty or could not be extracted"
            raise  # Re-raise the same exception to be caught by the caller
        finally:
            return parsed_content

