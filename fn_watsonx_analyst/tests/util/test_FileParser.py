"""Unit tests for File Parser module"""
import io
import os
import tempfile
from unittest.mock import patch
import pytest
import base64
from PIL import Image, ImageDraw
from pptx import Presentation
from docx import Document
import openpyxl
from fn_watsonx_analyst.util.FileParser import FileParser


class TestFileParser:
    """Test suite for the FileParser class."""

    @classmethod
    def setup_class(cls):
        """TBD"""
        cls.parser = FileParser()


    def read_pdf_from_data_folder(self, filename):
        """Helper function to read a PDF as bytes from the data folder."""
        file_path = os.path.join(os.path.dirname(__file__), '..', 'data', filename)
        with open(file_path, 'rb') as f:
            return f.read()

    @pytest.mark.livetest
    @patch("pdf2image.convert_from_bytes", return_value=[])  # Mock OCR image extraction
    def test_read_pdf_with_image_text(self, mock_convert):
        """Test PDF parsing with OCR fallback mocked."""
        parser = FileParser()

        # Use the helper function to read PDF bytes
        pdf_bytes = self.read_pdf_from_data_folder('test.pdf')
        result = parser.read_pdf(pdf_bytes)

        # Assert that expected text is in the result
        assert "when you need the model" in result

    def test_read_pdf(self):
        """Test for read_pdf with a minimal valid PDF byte stream."""
        # Base64 encoded minimal PDF with "Hello PDF" extractable text
        pdf_base64 = (
            "JVBERi0xLjQKJeLjz9MKMSAwIG9iago8PCAvVHlwZSAvQ2F0YWxvZyAvUGFnZXMgMiAwIFIg"
            "Pj4KZW5kb2JqCjIgMCBvYmoKPDwgL1R5cGUgL1BhZ2VzIC9Db3VudCAxIC9LaWRzIFsgMyAw"
            "IFIgXSA+PgplbmRvYmoKMyAwIG9iago8PCAvVHlwZSAvUGFnZSAvUGFyZW50IDIgMCBSIC9N"
            "ZWRpYUJveCBbMCAwIDYxMiA3OTJdIC9Db250ZW50cyA0IDAgUiAvUmVzb3VyY2VzIDw8IC9G"
            "b250IDw8IC9GMSA1IDAgUiA+PiA+PiA+PgplbmRvYmoKNCAwIG9iago8PCAvTGVuZ3RoIDQy"
            "ID4+CnN0cmVhbQpCVCAvRjEgMTIgVGYgMTAwIDcwMCBUZCAoSGVsbG8gUERGKSBUagpFVApl"
            "bmRzdHJlYW0KZW5kb2JqCjUgMCBvYmoKPDwgL1R5cGUgL0ZvbnQgL1N1YnR5cGUgL1R5cGUx"
            "IC9CYXNlRm9udCAvSGVsdmV0aWNhIC9FbmNvZGluZyAvV2luQW5zaUVuY29kaW5nID4+CmVu"
            "ZG9iagp4cmVmCjAgNgowMDAwMDAwMDAwIDY1NTM1IGYgCjAwMDAwMDAxMTggMDAwMDAgbiAK"
            "MDAwMDAwMDE4MSAwMDAwMCBuIAowMDAwMDAwMzAwIDAwMDAwIG4gCjAwMDAwMDAzODAgMDAw"
            "MDAgbiAKMDAwMDAwMDUwNCAwMDAwMCBuIAp0cmFpbGVyCjw8IC9TaXplIDYgL1Jvb3QgMSAw"
            "IFIgPj4Kc3RhcnR4cmVmCjYyNQolJUVPRgo="
        )
        pdf_bytes = base64.b64decode(pdf_base64)
        result = self.parser.read_pdf(pdf_bytes)
        assert "Hello PDF" in result

    def test_read_docx(self):
        """Test for read_docx"""
        doc = Document()
        doc.add_paragraph("This is a test DOCX.")
        with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp:
            doc.save(tmp.name)
            tmp.seek(0)
            data = tmp.read()

        result = self.parser.read_docx(data)
        assert "This is a test DOCX." in result

    def test_read_pptx(self):
        """Test for read_pptx"""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "This is a test PPTX."
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
            prs.save(tmp.name)
            tmp.seek(0)
            data = tmp.read()

        result = self.parser.read_pptx(data)
        assert "This is a test PPTX." in result

    def test_read_eml(self):
        """Test for read_eml"""
        eml_data = b"""\
From: sender@example.com
To: recipient@example.com
Subject: Test Email
MIME-Version: 1.0
Content-Type: text/plain

This is the plain text body of the email.
"""
        result = self.parser.read_eml(eml_data)
        assert "This is the plain text body of the email." in result

    def test_read_ics(self):
        """Test for read_ics and assert on the full result without using mock"""
        # Example of a valid ICS file content
        ics_data = """BEGIN:VCALENDAR
PRODID:-//Sample Calendar//EN
VERSION:2.0
BEGIN:VEVENT
SUMMARY:Meeting
DESCRIPTION:This is a test ICS event.
DTSTART:20250329T100000Z
DTEND:20250329T110000Z
END:VEVENT
END:VCALENDAR"""

        # Call the method on the actual `ics_data`
        result = self.parser.read_ics(ics_data)
        # Expected result based on the ICS content
        expected_result = "Event: Meeting, Start: 2025-03-29T10:00:00+00:00, End: 2025-03-29T11:00:00+00:00"

        # Assert the full result
        assert result == expected_result

    def test_read_xlsx(self):
        """Test for read_xlsx"""
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.append(["Column1", "Column2"])
        ws.append(["Value1", "Value2"])

        with tempfile.NamedTemporaryFile(suffix=".xlsx") as tmp:
            wb.save(tmp.name)
            tmp.seek(0)
            data = tmp.read()

        result = self.parser.read_xlsx(data)
        assert "Column1" in result
        assert "Value1" in result

    def test_read_ppt(self):
        """Test for read_ppt
        This test function creates a minimal .pptx file in memory, 
        replaces the actual conversion process with a dummy function, 
        and checks if the `read_ppt` method correctly reads the 
        file content without performing the conversion.
        """

        # Create a minimal in-memory .pptx file
        pptx_stream = io.BytesIO()
        prs = Presentation()
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Test Slide"
        prs.save(pptx_stream)
        pptx_bytes = pptx_stream.getvalue()

        # Define dummy_converter with 'self' since it'll be bound to the parser instance
        def dummy_converter(self, _: bytes) -> bytes:
            """
            The method takes bytes as input and returns bytes.
            _ is a placeholder for any byte data that might be passed to the function.
            """
            return pptx_bytes

        # Bind dummy_converter as method
        self.parser.convert_ppt_to_pptx = dummy_converter.__get__(self.parser)

        # Run the test
        result = self.parser.read_ppt(b"anything doesn't really matter")
        assert isinstance(result, str)
        assert "Test Slide" in result

    @patch("fn_watsonx_analyst.util.FileParser.image_to_string")
    def test_extract_text_from_images(self, mock_image_to_string):
        """Test for extract_text_from_images with real image data"""

        # Create a simple image with text using PIL
        img = Image.new('RGB', (200, 100), color=(255, 255, 255))
        draw = ImageDraw.Draw(img)
        draw.text((10, 40), "Detected text from image", fill=(0, 0, 0))

        # Now pass the real image object to the OCR method
        image_data = [img]  # List of real image objects (no mock)

        # Set the return value of the mocked function
        mock_image_to_string.return_value = "Detected text from image"

        # Call the method to extract text from images
        result = self.parser.extract_text_from_images(image_data)

        # Assert that the extracted text is as expected
        assert result == "Detected text from image"

    def test_format_parser(self):
        """Test for multi_format_parser (PDF)"""
        result = self.parser.multi_format_parser(b"%PDF-1.4 some pdf bytes", "test.pdf")
        assert isinstance(result, str)

    def test_format_parser_unknown(self):
        """Test for unsupported format"""
        result = self.parser.multi_format_parser(b"\xff\xfe\xfa\xfb", "test.abc")
        assert "Parsed content is empty or could not be extracted" in result
